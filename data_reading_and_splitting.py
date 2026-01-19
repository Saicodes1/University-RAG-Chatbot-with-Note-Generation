from langchain_experimental.text_splitter import SemanticChunker
from langchain_community.document_loaders import PDFPlumberLoader,UnstructuredPowerPointLoader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
from pptx import Presentation

embeddings=HuggingFaceEmbeddings(
    model_name="nomic-ai/nomic-embed-text-v1",
    model_kwargs={"trust_remote_code": True}
)
text_splitter=SemanticChunker(embeddings)

def loading_and_chunking_with_sentence(path):
    docs=[]
    ppt_docs=[]
    if (path.endswith("pdf")):
        loader= PDFPlumberLoader(path)
        docs=loader.load()
        #chunks=text_splitter.split_documents(docs)

    elif (path.endswith("pptx") or path.endswith("ppt") ):
        presentation = Presentation(path)
        num_slides=len(presentation.slides)
        ppt_docs=[]
        for i, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape,"text"):
                    if (shape.text in [str (j) for j in range(1,num_slides)] or ""):
                        continue
                    else:
                        x = Document(page_content=shape.text, metadata={"page_number": i})
                        ppt_docs.append(x)
                        #chunks=text_splitter.split_documents(ppt_docs)
    
    return docs, ppt_docs


def total_files():

    final_doc=[]
   # _,a=loading_and_chunking_with_sentence("Data\\CC-1.pptx")
    #for doc in a:
     #   doc.metadata.update({"Chapter_name":"Introduction to Compilers and Phases"})
    #chunks1=text_splitter.split_documents(a)
    #print("1st file is done")

    _,b=loading_and_chunking_with_sentence("Data\\CN-2.pptx")
    for doc in b:
        doc.metadata.update({"Chapter_name":"Chapter_2"})
    chunks2=text_splitter.split_documents(b)
    print("2nd file is done")

    _,c=loading_and_chunking_with_sentence("Data\\CN-3.pptx")
    for doc in c:
        doc.metadata.update({"Chapter_name":"Chapter_3"})
    chunks3=text_splitter.split_documents(c)
    print("3rd file is done")
    ''' d=loading_and_chunking_with_sentence("Data\\CC-4.pptx")
    print("4th file is done")
    e=loading_and_chunking_with_sentence("Data\\CC-5.pptx")
    print("5th file is done")
    f=loading_and_chunking_with_sentence("Data\\CC-6.pptx")
    print("6th file is done")
    g=loading_and_chunking_with_sentence("Data\\CC-7.pptx")
    print("7th file is done")
    h=loading_and_chunking_with_sentence("Data\\CC-8.pdf")
    print("8th file is done")
    i=loading_and_chunking_with_sentence("Data\\CC-9.pdf")
    print("9th file is done") '''
    print("all files are done")
    #final_doc.append(chunks1)
    final_doc.append(chunks2)
    final_doc.append(chunks3)
    '''
    final_doc.append(d)
    final_doc.append(e)
    final_doc.append(f)
    final_doc.append(g)
    final_doc.append(h)
    final_doc.append(i)
    '''
    return final_doc

